"""
Build Shah_Shahul_5_Presentation_032025.pptx
Professional deck for both technical and business audiences.
"""
import io
import json
import math
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patches as patches
from matplotlib.patches import FancyArrowPatch
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Brand palette
# ---------------------------------------------------------------------------
GREEN_DARK   = RGBColor(0x1B, 0x5E, 0x20)   # deep forest green
GREEN_MID    = RGBColor(0x2E, 0x7D, 0x32)   # rich green
GREEN_LIGHT  = RGBColor(0x66, 0xBB, 0x6A)   # leaf green
GOLD         = RGBColor(0xF9, 0xA8, 0x25)   # harvest gold
EARTH        = RGBColor(0x6D, 0x4C, 0x41)   # soil brown
CREAM        = RGBColor(0xF9, 0xFB, 0xE7)   # off-white
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x1A, 0x23, 0x1A)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_BLUE  = RGBColor(0x01, 0x57, 0x9B)   # water blue

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUT = Path(__file__).parent / "presentation_032025.pptx"
ASSETS = Path(__file__).parent.parent / "models"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def rgb_hex(r):
    return "#{:02X}{:02X}{:02X}".format(r[0], r[1], r[2])

def set_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_rect(slide, left, top, width, height, r, g, b, alpha_pct=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def add_fig(slide, fig, left, top, width, height):
    stream = fig_to_stream(fig)
    slide.shapes.add_picture(stream, left, top, width, height)


# ---------------------------------------------------------------------------
# Chart builders
# ---------------------------------------------------------------------------

def make_model_bar_chart():
    models = ["Ridge", "Hist GB", "LightGBM", "XGBoost", "Random\nForest"]
    r2     = [0.9130,  0.9128,    0.9128,     0.9127,    0.9099]
    colors = ["#2E7D32","#558B2F","#689F38","#8BC34A","#AED581"]

    fig, ax = plt.subplots(figsize=(7, 3.6), facecolor="#F9FBE7")
    ax.set_facecolor("#F9FBE7")
    bars = ax.barh(models[::-1], [v * 100 for v in r2[::-1]],
                   color=colors[::-1], height=0.55, edgecolor="white", linewidth=1.5)
    for bar, val in zip(bars, r2[::-1]):
        ax.text(bar.get_width() - 0.01, bar.get_y() + bar.get_height()/2,
                f"{val:.4f}", va="center", ha="right", fontsize=10,
                color="white", fontweight="bold")
    ax.set_xlim(90.5, 91.5)
    ax.set_xlabel("R² (%)", fontsize=10, color="#1A231A")
    ax.set_title("Model Comparison — R² Score", fontsize=12,
                 fontweight="bold", color="#1B5E20", pad=10)
    ax.tick_params(colors="#1A231A", labelsize=9)
    ax.spines[["top","right","bottom"]].set_visible(False)
    ax.xaxis.set_tick_params(length=0)
    fig.tight_layout(pad=1.0)
    return fig

def make_feature_importance_chart():
    features = ["rainfall_mm", "rainfall_anomaly", "rain_x_fertilizer",
                "harvest_rain_rate", "rain_x_irrigation", "water_stress",
                "gdd_proxy", "heat_stress"]
    corrs    = [0.764, 0.764, 0.668, 0.647, 0.594, 0.431, 0.312, 0.198]
    colors   = ["#1B5E20" if c > 0.6 else "#558B2F" if c > 0.4 else "#A5D6A7"
                for c in corrs]

    fig, ax = plt.subplots(figsize=(6.5, 3.5), facecolor="#F9FBE7")
    ax.set_facecolor("#F9FBE7")
    bars = ax.barh(features[::-1], corrs[::-1], color=colors[::-1],
                   height=0.55, edgecolor="white", linewidth=1)
    for bar, val in zip(bars, corrs[::-1]):
        ax.text(bar.get_width() + 0.01, bar.get_y() + bar.get_height()/2,
                f"{val:.3f}", va="center", ha="left", fontsize=9, color="#1A231A")
    ax.set_xlim(0, 0.9)
    ax.set_xlabel("|Pearson r|", fontsize=9, color="#1A231A")
    ax.set_title("Top Features by Correlation with Yield", fontsize=11,
                 fontweight="bold", color="#1B5E20", pad=8)
    ax.tick_params(colors="#1A231A", labelsize=8)
    ax.spines[["top","right","bottom"]].set_visible(False)
    fig.tight_layout(pad=1.0)
    return fig

def make_crop_yield_chart():
    crops = ["Maize", "Rice", "Soybean", "Wheat"]
    fao   = [9.93,   8.62,   2.96,      3.17]
    clrs  = ["#F9A825","#E53935","#7CB342","#0288D1"]

    fig, ax = plt.subplots(figsize=(5.5, 3.2), facecolor="#F9FBE7")
    ax.set_facecolor("#F9FBE7")
    x = np.arange(len(crops))
    bars = ax.bar(x, fao, color=clrs, width=0.55, edgecolor="white", linewidth=1.5)
    for bar, val in zip(bars, fao):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.1,
                f"{val} t/ha", ha="center", va="bottom", fontsize=9,
                fontweight="bold", color="#1A231A")
    ax.set_xticks(x)
    ax.set_xticklabels(crops, fontsize=10)
    ax.set_ylabel("FAO Benchmark Yield (t/ha)", fontsize=9, color="#1A231A")
    ax.set_title("FAO 2013 USA Reference Yields", fontsize=11,
                 fontweight="bold", color="#1B5E20", pad=8)
    ax.tick_params(colors="#1A231A")
    ax.spines[["top","right","left"]].set_visible(False)
    ax.yaxis.set_tick_params(length=0)
    ax.set_ylim(0, 12)
    fig.tight_layout(pad=1.0)
    return fig

def make_pipeline_diagram():
    fig, ax = plt.subplots(figsize=(10, 2.2), facecolor="#1B5E20")
    ax.set_facecolor("#1B5E20")
    ax.axis("off")

    stages = [
        ("Stage 1\nPreprocess", "#2E7D32"),
        ("Stage 2\nFeatures", "#388E3C"),
        ("Stage 3\nTrain", "#43A047"),
        ("Stage 4\nEvaluate", "#4CAF50"),
        ("FastAPI\nServe", "#F9A825"),
        ("Streamlit\nUI", "#66BB6A"),
    ]
    n = len(stages)
    w, h = 1.3, 0.8
    gap = 0.35
    total = n * w + (n - 1) * gap
    x0 = (10 - total) / 2

    for i, (label, color) in enumerate(stages):
        x = x0 + i * (w + gap)
        box = patches.FancyBboxPatch((x, 0.7), w, h,
                                      boxstyle="round,pad=0.06",
                                      facecolor=color, edgecolor="white",
                                      linewidth=1.5)
        ax.add_patch(box)
        ax.text(x + w/2, 0.7 + h/2, label, ha="center", va="center",
                fontsize=8.5, color="white", fontweight="bold",
                multialignment="center")
        if i < n - 1:
            ax.annotate("", xy=(x + w + gap, 1.1), xytext=(x + w + 0.02, 1.1),
                        arrowprops=dict(arrowstyle="->", color="#F9A825", lw=2))

    ax.set_xlim(0, 10)
    ax.set_ylim(0, 2.2)
    ax.set_title("AgroOpt — End-to-End MLOps Pipeline",
                 fontsize=11, color="white", fontweight="bold", pad=6)
    fig.tight_layout(pad=0.3)
    return fig

def make_architecture_diagram():
    fig, ax = plt.subplots(figsize=(8, 3.5), facecolor="#F9FBE7")
    ax.set_facecolor("#F9FBE7")
    ax.axis("off")

    def box(x, y, w, h, label, sub, color, tcolor="white"):
        rect = patches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.05",
                                       facecolor=color, edgecolor="white", linewidth=2)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h*0.65, label, ha="center", va="center",
                fontsize=9, color=tcolor, fontweight="bold")
        ax.text(x + w/2, y + h*0.28, sub, ha="center", va="center",
                fontsize=7.5, color=tcolor, alpha=0.85)

    def arrow(x1, y1, x2, y2):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", color="#2E7D32", lw=2))

    box(0.2, 1.8, 1.5, 0.9, "Farmer / User", "Browser / App", "#6D4C41")
    box(0.2, 0.5, 1.5, 0.9, "Streamlit UI", ":8501", "#2E7D32")
    box(3.0, 0.5, 1.8, 0.9, "FastAPI", ":8000  /predict\n/recommend  /optimize", "#1B5E20")
    box(6.0, 0.5, 1.5, 0.9, "Ridge Model", "best_model.pkl", "#01579B")
    box(6.0, 1.8, 1.5, 0.9, "fao_refs.json\nfeature_names.json", "models/", "#558B2F", "#FFFFFF")
    box(3.0, 1.8, 1.8, 0.9, "Docker Compose", "api + frontend\nservices", "#F9A825", "#1A231A")

    arrow(0.95, 1.8, 0.95, 1.4)
    arrow(1.7, 0.95, 3.0, 0.95)
    arrow(4.8, 0.95, 6.0, 0.95)
    arrow(6.75, 1.8, 6.75, 1.4)
    arrow(3.9, 1.8, 3.9, 1.4)

    ax.set_xlim(0, 8)
    ax.set_ylim(0.2, 3.0)
    ax.set_title("Deployment Architecture", fontsize=11,
                 fontweight="bold", color="#1B5E20", pad=8)
    fig.tight_layout(pad=0.5)
    return fig

def make_metrics_visual():
    fig, axes = plt.subplots(1, 4, figsize=(10, 2.2), facecolor="#1B5E20")
    fig.patch.set_facecolor("#1B5E20")
    metrics = [
        ("0.913", "R² Score", "#F9A825"),
        ("~0.5 t/ha", "Prediction Error", "#66BB6A"),
        ("666K", "Training Records", "#81D4FA"),
        ("106", "Automated Tests", "#FFCC80"),
    ]
    for ax, (val, label, color) in zip(axes, metrics):
        ax.set_facecolor("#2E7D32")
        ax.axis("off")
        rect = patches.FancyBboxPatch((0.05, 0.05), 0.9, 0.9,
                                       boxstyle="round,pad=0.04",
                                       facecolor="#2E7D32", edgecolor=color, linewidth=3,
                                       transform=ax.transAxes)
        ax.add_patch(rect)
        ax.text(0.5, 0.62, val, transform=ax.transAxes,
                ha="center", va="center", fontsize=20, fontweight="bold",
                color=color)
        ax.text(0.5, 0.28, label, transform=ax.transAxes,
                ha="center", va="center", fontsize=9, color="white")
    fig.suptitle("Key Performance Metrics", fontsize=12,
                 color="white", fontweight="bold", y=1.02)
    fig.tight_layout(pad=0.6)
    return fig

def make_data_overview_chart():
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(8, 3.2), facecolor="#F9FBE7")
    fig.patch.set_facecolor("#F9FBE7")

    # Pie: crop distribution (approx equal since synthetic)
    sizes = [25, 25, 25, 25]
    labels = ["Maize", "Rice", "Soybean", "Wheat"]
    clrs = ["#F9A825","#E53935","#7CB342","#0288D1"]
    wedges, texts, autotexts = ax1.pie(sizes, labels=labels, colors=clrs,
                                        autopct="%1.0f%%", startangle=90,
                                        wedgeprops=dict(edgecolor="white", linewidth=2),
                                        textprops=dict(fontsize=9, color="#1A231A"))
    for at in autotexts:
        at.set_color("white")
        at.set_fontweight("bold")
    ax1.set_title("Dataset by Crop", fontsize=10, fontweight="bold",
                  color="#1B5E20", pad=8)
    ax1.set_facecolor("#F9FBE7")

    # Bar: dataset split
    split_labels = ["Train", "Test"]
    split_vals   = [533195, 133299]
    split_clrs   = ["#2E7D32", "#A5D6A7"]
    ax2.set_facecolor("#F9FBE7")
    bars = ax2.bar(split_labels, [v/1000 for v in split_vals],
                   color=split_clrs, width=0.45, edgecolor="white", linewidth=1.5)
    for bar, val in zip(bars, split_vals):
        ax2.text(bar.get_x() + bar.get_width()/2,
                 bar.get_height() + 3, f"{val:,}", ha="center",
                 va="bottom", fontsize=9, fontweight="bold", color="#1A231A")
    ax2.set_ylabel("Records (thousands)", fontsize=9, color="#1A231A")
    ax2.set_title("80/20 Train-Test Split", fontsize=10, fontweight="bold",
                  color="#1B5E20", pad=8)
    ax2.spines[["top","right","left"]].set_visible(False)
    ax2.tick_params(colors="#1A231A")
    ax2.set_ylim(0, 650)

    fig.tight_layout(pad=1.0)
    return fig

def make_cicd_diagram():
    fig, ax = plt.subplots(figsize=(9, 3.0), facecolor="#F9FBE7")
    ax.set_facecolor("#F9FBE7")
    ax.axis("off")

    stages = [
        ("git push\nmain", "#6D4C41", "white"),
        ("Lint\n(flake8)", "#1B5E20", "white"),
        ("Test\n(pytest 106)", "#2E7D32", "white"),
        ("Docker\nBuild", "#388E3C", "white"),
        ("Docker\nPublish\nghcr.io", "#F9A825", "#1A231A"),
    ]
    n = len(stages)
    w, h = 1.35, 0.75
    gap = 0.3
    total = n * w + (n - 1) * gap
    x0 = (9 - total) / 2

    for i, (label, color, tcolor) in enumerate(stages):
        x = x0 + i * (w + gap)
        box = patches.FancyBboxPatch((x, 1.1), w, h,
                                      boxstyle="round,pad=0.06",
                                      facecolor=color, edgecolor="white", linewidth=2)
        ax.add_patch(box)
        ax.text(x + w/2, 1.1 + h/2, label, ha="center", va="center",
                fontsize=8.5, color=tcolor, fontweight="bold",
                multialignment="center")
        if i < n - 1:
            ax.annotate("", xy=(x + w + gap, 1.47), xytext=(x + w + 0.02, 1.47),
                        arrowprops=dict(arrowstyle="->", color="#2E7D32", lw=2.5))
        if i == 3:
            ax.annotate("", xy=(x + w + gap, 1.47), xytext=(x + w + 0.02, 1.47),
                        arrowprops=dict(arrowstyle="->", color="#F9A825",
                                        lw=2.5, linestyle="dashed"))

    ax.text(x0 + 3 * (w + gap) + w/2, 1.0,
            "tags only", ha="center", va="top", fontsize=7.5,
            color="#F9A825", style="italic")

    ax.set_xlim(0, 9)
    ax.set_ylim(0.5, 2.3)
    ax.set_title("CI/CD Pipeline — GitHub Actions", fontsize=11,
                 fontweight="bold", color="#1B5E20", pad=6)
    fig.tight_layout(pad=0.3)
    return fig


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, 0x1B, 0x5E, 0x20)

    # Dark green top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), 0x1B, 0x5E, 0x20)

    # Gold accent stripe
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(0.06), 0xF9, 0xA8, 0x25)

    # Decorative right panel
    add_rect(slide, Inches(8.8), 0, Inches(4.53), SLIDE_H, 0x15, 0x4A, 0x18)

    # Wheat stalk pattern block
    fig, ax = plt.subplots(figsize=(4.2, 7.5), facecolor="#154A18")
    ax.set_facecolor("#154A18")
    ax.axis("off")
    rng = np.random.default_rng(7)
    for _ in range(18):
        x = rng.uniform(0.1, 0.9)
        y0 = rng.uniform(0.0, 0.3)
        h  = rng.uniform(0.5, 0.85)
        color = rng.choice(["#66BB6A","#81C784","#A5D6A7","#F9A825"])
        alpha = rng.uniform(0.25, 0.55)
        ax.plot([x, x], [y0, y0 + h], color=color, lw=rng.uniform(1.2, 2.5),
                alpha=alpha, solid_capstyle="round")
        # head
        for j in range(5):
            dy = h * (0.55 + j * 0.08)
            ax.plot([x, x + rng.uniform(-0.06, 0.06)],
                    [y0 + dy, y0 + dy + 0.025],
                    color=color, lw=1.2, alpha=alpha * 1.2)
    # field horizon
    xs = np.linspace(0, 1, 200)
    ys = 0.15 + 0.03 * np.sin(xs * 12) + 0.02 * np.sin(xs * 30)
    ax.fill_between(xs, 0, ys, color="#6D4C41", alpha=0.6)
    fig.tight_layout(pad=0)
    buf = fig_to_stream(fig)
    slide.shapes.add_picture(buf, Inches(8.8), 0, Inches(4.53), SLIDE_H)

    # Sun / circle motif
    fig2, ax2 = plt.subplots(figsize=(2, 2), facecolor="#154A18")
    ax2.set_facecolor("#154A18")
    ax2.axis("off")
    circle = plt.Circle((0.5, 0.5), 0.42, color="#F9A825", zorder=2)
    ax2.add_patch(circle)
    inner  = plt.Circle((0.5, 0.5), 0.33, color="#154A18", zorder=3)
    ax2.add_patch(inner)
    for angle in np.linspace(0, 360, 16, endpoint=False):
        r = math.radians(angle)
        ax2.plot([0.5 + 0.44 * math.cos(r), 0.5 + 0.52 * math.cos(r)],
                 [0.5 + 0.44 * math.sin(r), 0.5 + 0.52 * math.sin(r)],
                 color="#F9A825", lw=2.5)
    ax2.set_xlim(0, 1); ax2.set_ylim(0, 1)
    fig2.tight_layout(pad=0)
    buf2 = fig_to_stream(fig2)
    slide.shapes.add_picture(buf2, Inches(9.3), Inches(0.4), Inches(1.5), Inches(1.5))

    # Title text
    add_text(slide, "AgroOpt",
             Inches(0.5), Inches(1.5), Inches(8.0), Inches(1.4),
             font_size=58, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(slide, "Agricultural Yield Prediction & Optimization",
             Inches(0.5), Inches(2.95), Inches(7.8), Inches(0.7),
             font_size=22, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

    add_text(slide, "An end-to-end MLOps platform for smarter farm decisions",
             Inches(0.5), Inches(3.65), Inches(7.8), Inches(0.5),
             font_size=15, bold=False, color=CREAM, align=PP_ALIGN.LEFT, italic=True)

    # Divider
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(5.0), Inches(0.04), 0xF9, 0xA8, 0x25)

    add_text(slide, "Shah Shahul  |  OpenClassrooms ML Engineer  |  March 2025",
             Inches(0.5), Inches(4.5), Inches(7.8), Inches(0.4),
             font_size=12, bold=False, color=CREAM, align=PP_ALIGN.LEFT)

    add_text(slide, "Crops: Maize  |  Rice  |  Soybean  |  Wheat",
             Inches(0.5), Inches(5.05), Inches(7.8), Inches(0.35),
             font_size=11, bold=False, color=GREEN_LIGHT, align=PP_ALIGN.LEFT)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0xF9, 0xFB, 0xE7)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)

    add_text(slide, "Agenda", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    items = [
        ("01", "Business Problem & Opportunity"),
        ("02", "Our Solution — AgroOpt Platform"),
        ("03", "Data & Methodology"),
        ("04", "Feature Engineering"),
        ("05", "Model Training & Results"),
        ("06", "Deployment Architecture"),
        ("07", "CI/CD & Testing"),
        ("08", "Business Value & ROI"),
        ("09", "Recommendations & Next Steps"),
    ]
    cols = 3
    col_w = Inches(4.2)
    row_h = Inches(1.05)
    for i, (num, text) in enumerate(items):
        col = i % cols
        row = i // cols
        lft = Inches(0.35) + col * col_w
        top = Inches(1.4) + row * row_h
        add_rect(slide, lft, top, Inches(0.55), Inches(0.55), 0x2E, 0x7D, 0x32)
        add_text(slide, num, lft, top, Inches(0.55), Inches(0.55),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, text, lft + Inches(0.62), top + Inches(0.04),
                 Inches(3.5), Inches(0.5),
                 font_size=13, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0xF9, 0xFB, 0xE7)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "The Business Problem", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.7), font_size=30, bold=True, color=WHITE)

    # Three pain-point cards
    cards = [
        ("What to Plant?",
         "Which crop will perform best given local soil, climate & region?",
         "#1B5E20", "#F9A825"),
        ("How Much Will I Yield?",
         "What is a realistic yield forecast before the season starts?",
         "#01579B", "#81D4FA"),
        ("How Can I Improve?",
         "Which management changes (fertilizer, irrigation, timing) maximise output?",
         "#6D4C41", "#FFCC80"),
    ]
    for i, (title, body, bg, accent) in enumerate(cards):
        lft = Inches(0.4 + i * 4.3)
        add_rect(slide, lft, Inches(1.35), Inches(4.0), Inches(2.5),
                 *[int(bg.lstrip("#")[j:j+2], 16) for j in (0, 2, 4)])
        add_rect(slide, lft, Inches(1.35), Inches(4.0), Inches(0.12),
                 *[int(accent.lstrip("#")[j:j+2], 16) for j in (0, 2, 4)])
        add_text(slide, title, lft + Inches(0.15), Inches(1.6),
                 Inches(3.7), Inches(0.55),
                 font_size=14, bold=True, color=WHITE)
        add_text(slide, body, lft + Inches(0.15), Inches(2.2),
                 Inches(3.7), Inches(0.9),
                 font_size=11.5, bold=False, color=CREAM)

    # Stats
    add_rect(slide, Inches(0.4), Inches(4.05), Inches(12.4), Inches(0.05),
             0xF9, 0xA8, 0x25)
    stats = [
        "$200B+", "US crop production annual value",
        "70%",    "of global freshwater used in agriculture",
        "5-17%",  "yield improvement from data-driven decisions",
    ]
    for i in range(0, len(stats), 2):
        lft = Inches(0.5 + (i//2) * 4.4)
        add_text(slide, stats[i], lft, Inches(4.2), Inches(2.5), Inches(0.7),
                 font_size=28, bold=True, color=GREEN_MID)
        add_text(slide, stats[i+1], lft, Inches(4.9), Inches(4.0), Inches(0.5),
                 font_size=11, bold=False, color=DARK_TEXT)


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x15, 0x4A, 0x18)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "The AgroOpt Solution", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.7), font_size=30, bold=True, color=WHITE)

    fig = make_pipeline_diagram()
    add_fig(slide, fig, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.5))

    bullets = [
        ("Predict", "Forecast yield in hg/ha for any of 4 crops given farm conditions"),
        ("Recommend", "Rank all 4 crops and show water/heat stress indices"),
        ("Optimize", "Grid-search fertilizer, irrigation & harvest timing to maximise yield"),
        ("Track", "Every experiment logged to MLflow for full reproducibility"),
    ]
    for i, (label, desc) in enumerate(bullets):
        lft = Inches(0.4 + i * 3.2)
        add_rect(slide, lft, Inches(4.05), Inches(3.0), Inches(2.8),
                 0x2E, 0x7D, 0x32)
        add_rect(slide, lft, Inches(4.05), Inches(3.0), Inches(0.1),
                 0xF9, 0xA8, 0x25)
        add_text(slide, label, lft + Inches(0.12), Inches(4.25),
                 Inches(2.8), Inches(0.45),
                 font_size=14, bold=True, color=GOLD)
        add_text(slide, desc, lft + Inches(0.12), Inches(4.75),
                 Inches(2.78), Inches(1.8),
                 font_size=11, bold=False, color=CREAM)


def slide_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0xF9, 0xFB, 0xE7)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "Data & Methodology", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.7), font_size=30, bold=True, color=WHITE)

    # Left: data overview chart
    fig = make_data_overview_chart()
    add_fig(slide, fig, Inches(0.4), Inches(1.3), Inches(7.0), Inches(3.2))

    # Right: data source table-style
    sources = [
        ("crop_yield.csv", "666,494 farm records (synthetic USA)"),
        ("FAO yield.csv",  "Historical USA crop yields (hg/ha)"),
        ("FAO rainfall.csv","USA annual rainfall (mm)"),
        ("FAO temp.csv",    "USA temperature (52-state mean)"),
    ]
    add_text(slide, "Data Sources", Inches(7.8), Inches(1.3), Inches(5.0), Inches(0.4),
             font_size=13, bold=True, color=GREEN_DARK)
    for i, (src, desc) in enumerate(sources):
        top = Inches(1.75 + i * 0.72)
        add_rect(slide, Inches(7.8), top, Inches(5.0), Inches(0.62), 0xE8, 0xF5, 0xE9)
        add_rect(slide, Inches(7.8), top, Inches(0.08), Inches(0.62), 0x2E, 0x7D, 0x32)
        add_text(slide, src, Inches(8.0), top + Inches(0.03),
                 Inches(4.7), Inches(0.28), font_size=10.5, bold=True, color=GREEN_DARK)
        add_text(slide, desc, Inches(8.0), top + Inches(0.3),
                 Inches(4.7), Inches(0.28), font_size=9.5, bold=False, color=DARK_TEXT)

    # Bottom stat strip
    add_rect(slide, 0, Inches(4.65), SLIDE_W, Inches(0.06), 0x2E, 0x7D, 0x32)
    kpis = [("666,494", "Total Records"), ("4", "Crops"),
            ("4", "US Regions"), ("0", "Missing Values")]
    for i, (val, lbl) in enumerate(kpis):
        lft = Inches(0.8 + i * 3.1)
        add_text(slide, val, lft, Inches(4.85), Inches(2.5), Inches(0.65),
                 font_size=26, bold=True, color=GREEN_MID)
        add_text(slide, lbl, lft, Inches(5.5), Inches(2.5), Inches(0.4),
                 font_size=11, bold=False, color=DARK_TEXT)

    # Crop yield chart
    fig2 = make_crop_yield_chart()
    add_fig(slide, fig2, Inches(7.7), Inches(4.6), Inches(5.3), Inches(2.7))


def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0xF9, 0xFB, 0xE7)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "Feature Engineering — 36 Predictive Features",
             Inches(0.5), Inches(0.2), Inches(12.5), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    # Feature correlation chart
    fig = make_feature_importance_chart()
    add_fig(slide, fig, Inches(0.4), Inches(1.3), Inches(6.8), Inches(3.5))

    # Feature category boxes
    categories = [
        ("Domain Physics", "Water stress, heat stress,\nGDD proxy, aridity index,\nsoil quality score", "#1B5E20"),
        ("Interaction Terms", "rainfall x fertilizer\nrainfall x irrigation\nGDD x irrigation", "#01579B"),
        ("Anomaly Features", "Rainfall anomaly\nvs FAO benchmark\nTemp anomaly", "#6D4C41"),
        ("One-Hot Encoded", "Crop (4), Region (4)\nSoil type (6)\nWeather (3)", "#388E3C"),
    ]
    for i, (cat, detail, color) in enumerate(categories):
        col = i % 2
        row = i // 2
        lft = Inches(7.5 + col * 2.9)
        top = Inches(1.35 + row * 2.0)
        add_rect(slide, lft, top, Inches(2.7), Inches(1.75),
                 *[int(color.lstrip("#")[j:j+2], 16) for j in (0, 2, 4)])
        add_text(slide, cat, lft + Inches(0.12), top + Inches(0.12),
                 Inches(2.5), Inches(0.4), font_size=11, bold=True, color=GOLD)
        add_text(slide, detail, lft + Inches(0.12), top + Inches(0.55),
                 Inches(2.5), Inches(1.1), font_size=9.5, bold=False, color=CREAM)

    add_text(slide,
             "Rainfall dominates yield (r=0.764) — consistent with US agricultural science",
             Inches(0.4), Inches(5.05), Inches(9.0), Inches(0.4),
             font_size=11, bold=True, color=GREEN_DARK, italic=True)


def slide_models(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0xF9, 0xFB, 0xE7)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "Model Training & Results", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.7), font_size=30, bold=True, color=WHITE)

    # Bar chart
    fig = make_model_bar_chart()
    add_fig(slide, fig, Inches(0.4), Inches(1.25), Inches(7.2), Inches(3.8))

    # Winner callout
    add_rect(slide, Inches(7.9), Inches(1.25), Inches(4.9), Inches(3.8),
             0x1B, 0x5E, 0x20)
    add_rect(slide, Inches(7.9), Inches(1.25), Inches(4.9), Inches(0.12),
             0xF9, 0xA8, 0x25)
    add_text(slide, "Winner: Ridge Regression",
             Inches(8.05), Inches(1.5), Inches(4.6), Inches(0.5),
             font_size=14, bold=True, color=GOLD)

    winner_stats = [
        ("R²",        "0.9130"),
        ("RMSE",      "4,989 hg/ha"),
        ("MAE",       "3,976 hg/ha"),
        ("Train time","1.4 seconds"),
    ]
    for i, (k, v) in enumerate(winner_stats):
        top = Inches(2.15 + i * 0.6)
        add_text(slide, k + ":", Inches(8.1), top, Inches(1.5), Inches(0.45),
                 font_size=11, bold=True, color=CREAM)
        add_text(slide, v, Inches(9.4), top, Inches(3.2), Inches(0.45),
                 font_size=11, bold=False, color=WHITE)

    add_text(slide,
             "Why Ridge wins: the synthetic dataset has predominantly\n"
             "linear relationships. All 5 models cluster at R2 ~0.91,\n"
             "confirming Ridge is the leanest choice.",
             Inches(8.05), Inches(4.45), Inches(4.6), Inches(0.55),
             font_size=9.5, bold=False, color=CREAM, italic=True)

    add_text(slide, "MLflow tracks every run: parameters, R2, RMSE, MAE, MAPE, CV-RMSE",
             Inches(0.4), Inches(5.2), Inches(13.0), Inches(0.4),
             font_size=11, bold=False, color=DARK_TEXT, italic=True)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0xF9, 0xFB, 0xE7)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "Deployment Architecture", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.7), font_size=30, bold=True, color=WHITE)

    fig = make_architecture_diagram()
    add_fig(slide, fig, Inches(0.4), Inches(1.3), Inches(8.2), Inches(3.8))

    # Endpoint table
    endpoints = [
        ("GET",  "/health",    "Model status, version, feature count"),
        ("POST", "/predict",   "Yield forecast for a specific crop"),
        ("POST", "/recommend", "Ranked crop list with stress indices"),
        ("POST", "/optimize",  "Best mgmt settings to maximise yield"),
    ]
    add_text(slide, "REST Endpoints", Inches(9.0), Inches(1.3), Inches(4.0), Inches(0.4),
             font_size=13, bold=True, color=GREEN_DARK)
    for i, (method, ep, desc) in enumerate(endpoints):
        top = Inches(1.8 + i * 0.82)
        r, g, b = (0x2E, 0x7D, 0x32) if method == "GET" else (0x01, 0x57, 0x9B)
        add_rect(slide, Inches(9.0), top, Inches(0.65), Inches(0.5), r, g, b)
        add_text(slide, method, Inches(9.0), top, Inches(0.65), Inches(0.5),
                 font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, ep, Inches(9.7), top + Inches(0.02), Inches(3.1), Inches(0.25),
                 font_size=10, bold=True, color=GREEN_DARK)
        add_text(slide, desc, Inches(9.7), top + Inches(0.27), Inches(3.1), Inches(0.25),
                 font_size=9, bold=False, color=DARK_TEXT)

    add_text(slide, "Docker Compose: two services (api + frontend), non-root user, health checks",
             Inches(0.4), Inches(5.25), Inches(12.5), Inches(0.35),
             font_size=10.5, bold=False, color=DARK_TEXT, italic=True)


def slide_cicd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0xF9, 0xFB, 0xE7)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "CI/CD Pipeline & Testing", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.7), font_size=30, bold=True, color=WHITE)

    fig = make_cicd_diagram()
    add_fig(slide, fig, Inches(0.5), Inches(1.3), Inches(12.0), Inches(3.0))

    test_cards = [
        ("30 Tests", "Schema Validation", "Pydantic field bounds,\nenums, boundary values"),
        ("40 Tests", "Engine Unit Tests", "FarmConditions, feature\nvectors, predictions"),
        ("36 Tests", "API Integration", "All 4 endpoints,\nvalid & invalid inputs"),
    ]
    for i, (count, title, detail) in enumerate(test_cards):
        lft = Inches(0.5 + i * 4.2)
        add_rect(slide, lft, Inches(4.5), Inches(3.9), Inches(2.5), 0x1B, 0x5E, 0x20)
        add_rect(slide, lft, Inches(4.5), Inches(3.9), Inches(0.1), 0xF9, 0xA8, 0x25)
        add_text(slide, count, lft + Inches(0.15), Inches(4.7),
                 Inches(3.6), Inches(0.55), font_size=22, bold=True, color=GOLD)
        add_text(slide, title, lft + Inches(0.15), Inches(5.3),
                 Inches(3.6), Inches(0.4), font_size=12, bold=True, color=WHITE)
        add_text(slide, detail, lft + Inches(0.15), Inches(5.75),
                 Inches(3.6), Inches(0.9), font_size=10, bold=False, color=CREAM)


def slide_business_value(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x15, 0x4A, 0x18)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "Business Value", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.7), font_size=30, bold=True, color=WHITE)

    fig = make_metrics_visual()
    add_fig(slide, fig, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.2))

    value_pts = [
        ("Prediction Accuracy",
         "5% error on Maize (highest-volume crop). Reliable enough for planting decisions."),
        ("Crop Recommendation",
         "Instant ranking of 4 crops with water & heat stress indices for any location."),
        ("Yield Optimisation",
         "Grid search over fertilizer, irrigation and harvest timing returns yield-maximising plan."),
        ("Full Reproducibility",
         "MLflow logs every experiment. Any run can be reproduced from tracked parameters."),
        ("One-Command Deploy",
         "Docker Compose brings up the entire platform in seconds on any machine."),
        ("Continuous Confidence",
         "106 automated tests + CI pipeline catch regressions before they reach production."),
    ]
    cols = 2
    col_w = Inches(6.2)
    for i, (title, body) in enumerate(value_pts):
        col = i % cols
        row = i // cols
        lft = Inches(0.35 + col * col_w)
        top = Inches(3.75 + row * 1.05)
        add_rect(slide, lft, top, Inches(0.08), Inches(0.75), 0xF9, 0xA8, 0x25)
        add_text(slide, title, lft + Inches(0.18), top,
                 Inches(5.8), Inches(0.35), font_size=11, bold=True, color=GOLD)
        add_text(slide, body, lft + Inches(0.18), top + Inches(0.35),
                 Inches(5.8), Inches(0.45), font_size=9.5, bold=False, color=CREAM)


def slide_recommendations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0xF9, 0xFB, 0xE7)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "Recommendations & Roadmap", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.7), font_size=30, bold=True, color=WHITE)

    horizons = [
        ("Short-Term", "#1B5E20", [
            "Replace synthetic data with real USDA/NASS records",
            "Add prediction intervals (quantile regression)",
            "Reintroduce pesticide feature with real data",
        ]),
        ("Medium-Term", "#01579B", [
            "County-level granularity (beyond 4 regions)",
            "NOAA 90-day weather forecast integration",
            "Automated retraining + drift detection",
        ]),
        ("Long-Term", "#6D4C41", [
            "Satellite imagery (NDVI) as soil health proxy",
            "Shift objective: max yield -> max profit",
            "Multi-country expansion via FAO global data",
        ]),
    ]
    for i, (horizon, color, items) in enumerate(horizons):
        lft = Inches(0.4 + i * 4.3)
        add_rect(slide, lft, Inches(1.3), Inches(4.1), Inches(5.5),
                 *[int(color.lstrip("#")[j:j+2], 16) for j in (0, 2, 4)])
        add_rect(slide, lft, Inches(1.3), Inches(4.1), Inches(0.12),
                 0xF9, 0xA8, 0x25)
        add_text(slide, horizon, lft + Inches(0.15), Inches(1.5),
                 Inches(3.8), Inches(0.45), font_size=15, bold=True, color=GOLD)
        for j, item in enumerate(items):
            top = Inches(2.15 + j * 1.45)
            add_rect(slide, lft + Inches(0.15), top, Inches(0.35), Inches(0.35),
                     0xF9, 0xA8, 0x25)
            add_text(slide, str(j + 1), lft + Inches(0.15), top,
                     Inches(0.35), Inches(0.35), font_size=10, bold=True,
                     color=DARK_TEXT, align=PP_ALIGN.CENTER)
            add_text(slide, item, lft + Inches(0.6), top,
                     Inches(3.3), Inches(0.8), font_size=10.5,
                     bold=False, color=CREAM)


def slide_mlflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0xF9, 0xFB, 0xE7)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), 0x1B, 0x5E, 0x20)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), 0xF9, 0xA8, 0x25)
    add_text(slide, "MLflow Experiment Tracking", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.7), font_size=30, bold=True, color=WHITE)

    screenshot = Path(__file__).parent / "mlflow_screenshot_032025.png"
    slide.shapes.add_picture(str(screenshot),
                             Inches(0.5), Inches(1.25),
                             Inches(12.3), Inches(4.75))

    add_text(slide,
             "5 models tracked — Ridge, Hist GB, LightGBM, XGBoost, Random Forest  |  Metrics: R2, RMSE, MAE, MAPE, CV-RMSE",
             Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.45),
             font_size=10.5, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.CENTER, italic=True)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0x1B, 0x5E, 0x20)

    # Gold stripe
    add_rect(slide, 0, Inches(3.55), SLIDE_W, Inches(0.08), 0xF9, 0xA8, 0x25)

    # Decorative field left panel
    fig, ax = plt.subplots(figsize=(3.5, 7.5), facecolor="#154A18")
    ax.set_facecolor("#154A18")
    ax.axis("off")
    rng = np.random.default_rng(99)
    for _ in range(22):
        x  = rng.uniform(0.05, 0.95)
        y0 = rng.uniform(0.0, 0.25)
        h  = rng.uniform(0.4, 0.85)
        c  = rng.choice(["#66BB6A","#A5D6A7","#F9A825","#81C784"])
        a  = rng.uniform(0.3, 0.65)
        ax.plot([x, x], [y0, y0 + h], color=c, lw=rng.uniform(1.5, 3), alpha=a)
        for j in range(6):
            dy = h * (0.45 + j * 0.09)
            ax.plot([x, x + rng.uniform(-0.07, 0.07)],
                    [y0 + dy, y0 + dy + 0.03], color=c, lw=1.3, alpha=a * 1.3)
    xs = np.linspace(0, 1, 200)
    ys = 0.12 + 0.025 * np.sin(xs * 15)
    ax.fill_between(xs, 0, ys, color="#6D4C41", alpha=0.7)
    fig.tight_layout(pad=0)
    buf = fig_to_stream(fig)
    slide.shapes.add_picture(buf, 0, 0, Inches(3.3), SLIDE_H)

    add_text(slide, "Thank You", Inches(3.8), Inches(1.0),
             Inches(9.0), Inches(1.2), font_size=52, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)

    add_text(slide, "AgroOpt — Smarter Decisions, Better Harvests",
             Inches(3.8), Inches(2.35), Inches(9.0), Inches(0.6),
             font_size=18, bold=False, color=GOLD, align=PP_ALIGN.LEFT, italic=True)

    add_text(slide,
             "Shah Shahul\nOpenClassrooms — Machine Learning Engineer\nMarch 2025",
             Inches(3.8), Inches(4.0), Inches(9.0), Inches(1.2),
             font_size=13, bold=False, color=CREAM, align=PP_ALIGN.LEFT)

    add_text(slide,
             "Questions?",
             Inches(3.8), Inches(5.5), Inches(4.0), Inches(0.7),
             font_size=26, bold=True, color=GREEN_LIGHT, align=PP_ALIGN.LEFT)

    # GitHub repo note
    add_text(slide,
             "github.com/shah-data-scientist/Agro_Opt",
             Inches(3.8), Inches(6.3), Inches(8.0), Inches(0.45),
             font_size=11, bold=False, color=GREEN_LIGHT,
             align=PP_ALIGN.LEFT, italic=True)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_title(prs);          print("  1/10  Title")
    slide_agenda(prs);         print("  2/10  Agenda")
    slide_problem(prs);        print("  3/10  Problem")
    slide_solution(prs);       print("  4/10  Solution")
    slide_data(prs);           print("  5/10  Data")
    slide_features(prs);       print("  6/10  Features")
    slide_models(prs);         print("  7/10  Models")
    slide_mlflow(prs);         print("  8/13  MLflow Screenshot")
    slide_architecture(prs);   print("  9/13  Architecture")
    slide_cicd(prs);           print("  9/10  CI/CD")
    slide_business_value(prs); print(" 10/10  Business Value")
    slide_recommendations(prs);print(" 11/11  Recommendations")
    slide_closing(prs);        print(" 12/12  Closing")

    prs.save(OUT)
    print(f"Saved -> {OUT}")


if __name__ == "__main__":
    main()
